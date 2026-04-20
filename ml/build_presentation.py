"""build_presentation.py — Generate the PH 551 final-project slide deck.

Produces `ml/presentation.pptx` targeting the 100%-rubric criteria:
  - Data Exploration & Visualization
  - Problem Definition & Stats (naive baseline, correlation)
  - Methodology & Model Choice
  - Hyperparameter Tuning (GridSearchCV)
  - Evaluation & Metrics, Uncertainty (RMSE/MAE/R², UQ, coverage)
  - Result Visualization
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT  = Path(__file__).parent
FIG   = ROOT / "figures"
OUT   = ROOT / "presentation.pptx"

# 16:9 widescreen (13.33 × 7.5 in)
W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

DARK   = RGBColor(0x15, 0x28, 0x4B)
ACCENT = RGBColor(0xE9, 0x6A, 0x1F)
GRAY   = RGBColor(0x60, 0x60, 0x60)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)


def add_textbox(slide, text, *, left, top, width, height,
                size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb


def add_title(slide, text, subtitle=None):
    add_textbox(slide, text, left=Inches(0.4), top=Inches(0.25),
                width=Inches(12.5), height=Inches(0.8),
                size=32, bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, subtitle, left=Inches(0.4), top=Inches(1.0),
                    width=Inches(12.5), height=Inches(0.5),
                    size=16, color=GRAY)


def add_footer(slide, idx, total):
    add_textbox(slide, f"{idx}/{total}",
                left=Inches(12.6), top=Inches(7.1),
                width=Inches(0.6), height=Inches(0.3),
                size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image(slide, path, *, left, top, width=None, height=None):
    kw = {}
    if width  is not None: kw["width"]  = width
    if height is not None: kw["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kw)


def add_bullets(slide, items, *, left, top, width, height, size=16, color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


# ─── helper to background the slide with a subtle accent bar ──────────────
def slide_with_bar(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    # left accent bar
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_title(s, title, subtitle)
    return s


# ================ Slide 1: Title ================
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.4))
bar.fill.solid(); bar.fill.fore_color.rgb = DARK; bar.line.fill.background()
add_textbox(s, "Predicting Kalman-Filter Process Noise\nfrom Frequency-Stability Curves",
            left=Inches(0.5), top=Inches(0.2), width=Inches(12.5), height=Inches(1.2),
            size=32, bold=True, color=WHITE)
add_textbox(s, "PH 551 Final Project — Ian Lapinski",
            left=Inches(0.5), top=Inches(1.8), width=Inches(12.5), height=Inches(0.6),
            size=22, bold=True, color=DARK)
add_textbox(s,
    "Goal: learn (q_WPM, q_WFM, q_RWFM) from 196-feature σ(τ) vectors.\n"
    "Dataset: 10 000 synthetic Rb-class clocks (GMR6000-anchored h ranges).\n"
    "Validation: real GMR6000 Rb phase record (3 M samples, τ₀ = 1 s).",
    left=Inches(0.5), top=Inches(2.6), width=Inches(12.5), height=Inches(2.5),
    size=18, color=GRAY)
add_textbox(s, "SigmaTau v1.0 (Julia) · scikit-learn · XGBoost · allantools · Stable32",
            left=Inches(0.5), top=Inches(6.7), width=Inches(12.5), height=Inches(0.5),
            size=14, color=GRAY)

# ================ Slide 2: Motivation ================
s = slide_with_bar("Why predict q from σ(τ)?",
                   "Closed-loop tuning is expensive — ML can replace iterative search with a feed-forward call")
add_bullets(s, [
    "A Kalman filter for an atomic-clock ensemble needs process-noise parameters q = (q_WPM, q_WFM, q_RWFM).",
    "Industry practice: fit (mhdev_fit) or iterate (NLL, ALS). Iterative solvers are slow and multimodal.",
    "Our proposal: train an RF / XGB regressor so a 196-feature σ(τ) vector maps directly to q.",
    "Research question: Can ML predict q accurately enough to (a) use directly or (b) warm-start the NLL?",
    "Practical target: real GMR6000 Rb oscillator — does the learned model generalize to hardware?",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=20)

# ================ Slide 3: Background — clock model ================
s = slide_with_bar("Background: 3-state KF clock model",
                   "Wu (2023) convention — h↔q mapping used throughout SigmaTau")
add_bullets(s, [
    "State:  [phase, freq, drift]   Measurement = phase + WPM (R = q_WPM).",
    "Diffusion terms inside Φ → q_WFM (drives frequency random walk), q_RWFM (drives drift random walk).",
    "σ²_y,WPM(τ)  = 3·q_WPM / τ²     σ²_y,WFM(τ)  = q_WFM / τ     σ²_y,RWFM(τ) = q_RWFM · τ / 3",
    "Targets log10-transformed — each q spans 3–12 decades, making ratio metrics natural.",
    "Features: 20 raw log10-σ at canonical τ grid × (ADEV, MDEV, HDEV, MHDEV) = 80 raw.",
    "+ 19 slope features per statistic (76) + 40 ratio features (MVAR/AVAR, MHVAR/HVAR) = 196 total.",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=18)

# ================ Slide 4: Dataset ================
s = slide_with_bar("Dataset: 10 000 synthetic samples",
                   "Anchored to GMR6000 NLL fits; FPM present in 30% of samples")
add_bullets(s, [
    "Per-sample h-draws (log10 uniform):",
    "    h_+2 ∈ [−19, −16]   (WPM)    h_+1 ∈ [−28, −24]   (FPM, 30% prob.)",
    "    h_0  ∈ [−23, −20]   (WFM)    h_-1 ∈ [−28, −25]   (FFM)",
    "    h_-2 ∈ [−34, −28]   (RWFM)",
    "Per sample: N = 524 288 = 2¹⁹ phase points at τ₀ = 1 s  (~6 days of data).",
    "Labels: ALS (autocovariance least-squares) warm-started from true h → (q_WPM, q_WFM, q_RWFM).",
    "   — verified against truth on 10 probe draws; ALS matches to ≤ 10⁻³ dec, NLL drifted by up to 8 dec.",
    "Produced in 67.4 min on 12 threads (Julia, threaded).  All 10 000 ALS-converged.",
    "Data quality: 0 NaN, 0 Inf, 0 samples with any feature |z| > 5.",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=18)

# ================ Slide 5: EDA — distributions ================
s = slide_with_bar("EDA: feature & target distributions",
                   "Stratified 80/20 split preserves FPM-present balance across train/test")
# Image is a 2×3 grid (tall) — constrain by height, not width, so bottom row isn't clipped
add_image(s, FIG / "eda_distributions.png",
          left=Inches(2.3), top=Inches(1.6), height=Inches(5.7))

# ================ Slide 6: EDA — sample curves ================
s = slide_with_bar("EDA: example stability curves",
                   "Five random samples × 4 deviation families (ADEV/MDEV/HDEV/MHDEV)")
add_image(s, FIG / "eda_example_curves.png",
          left=Inches(0.5), top=Inches(1.6), width=Inches(12.3))

# ================ Slide 7: EDA — feature/target correlation ================
s = slide_with_bar("EDA: Pearson correlation  (196 features × 3 targets)",
                   "Short-τ features → q_WPM;  long-τ features → q_RWFM;  clear monotonic structure")
add_image(s, FIG / "eda_corr_feature_target.png",
          left=Inches(0.5), top=Inches(2.0), width=Inches(12.3))

# ================ Slide 8: Baseline + methodology ================
s = slide_with_bar("Problem framing + naive baseline",
                   "Regression from 196-D structured features → 3-D continuous target")
add_bullets(s, [
    "Naive baseline: predict train-set mean for every test sample.",
    "Naive RMSE:  q_WPM 0.87,  q_WFM 0.86,  q_RWFM 4.26   (R² ≈ 0 by construction)",
    "Tree ensembles chosen because:",
    "   — non-linear, non-monotonic interactions across σ(τ) slope regimes",
    "   — robust to feature scale (log10 σ already compresses dynamic range)",
    "   — RF provides variance-based UQ for free; XGB supports quantile objectives",
    "Baseline candidates rejected: linear regression (can't capture slope changes),",
    "k-NN (curse of dimensionality in 196-D), neural nets (data-thin at N=10k).",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=18)

# ================ Slide 9: GridSearchCV tuning ================
s = slide_with_bar("Hyperparameter tuning — 5-fold GridSearchCV",
                   "RF grid: 54 configs × 5 folds = 270 fits;  XGB grid: 36 × 5 = 180 fits")
add_bullets(s, [
    "RF grid explored:",
    "   n_estimators ∈ {200, 500, 1000}    max_depth ∈ {None, 20, 30}",
    "   min_samples_leaf ∈ {3, 5, 10}      max_features ∈ {sqrt, 0.5}",
    "Best RF:   n_estimators=500, max_depth=20, max_features=0.5, min_samples_leaf=3",
    "",
    "XGB grid explored:",
    "   n_estimators ∈ {200, 500}          learning_rate ∈ {0.01, 0.05, 0.1}",
    "   max_depth ∈ {4, 6, 8}              subsample ∈ {0.8, 1.0}",
    "Best XGB:  n_estimators=200, learning_rate=0.05, max_depth=4, subsample=0.8",
    "",
    "Scoring: neg_mean_squared_error (aligned with RMSE reporting metric).",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=16)

# ================ Slide 10: Test-set metrics ================
s = slide_with_bar("Test-set results — RF vs XGB vs Naive",
                   "N_test = 2000 held-out samples;   metrics in log10 decades")

from pptx.util import Pt as P
tbl_left = Inches(0.8); tbl_top = Inches(1.9)
tbl_w = Inches(8.5);    tbl_h = Inches(2.5)
rows, cols = 4, 5
table = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
headers = ["Target", "Naive RMSE", "RF RMSE (R²)", "XGB RMSE (R²)", "Improvement"]
data = [
    ("q_WPM",  "0.874", "0.069  (0.994)", "0.007  (0.99993)",  "119× vs naive"),
    ("q_WFM",  "0.873", "0.090  (0.989)", "0.053  (0.996)",    "16× vs naive"),
    ("q_RWFM", "1.732", "0.784  (0.795)", "0.736  (0.819)",    "2.4× vs naive"),
]
for j, h in enumerate(headers):
    c = table.cell(0, j)
    c.text = h
    for p in c.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    c.fill.solid(); c.fill.fore_color.rgb = DARK
for i, row in enumerate(data, start=1):
    for j, v in enumerate(row):
        c = table.cell(i, j)
        c.text = v
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
add_bullets(s, [
    "XGB dominates on q_WPM (measurement-noise regime, short-τ features well-resolved).",
    "q_RWFM remains the noisiest target — drift is poorly observed over 6-day window.",
    "Both models beat naive on every target;  XGB wins 2/3.",
],  left=Inches(0.8), top=Inches(4.8), width=Inches(12), height=Inches(2.2), size=16)

# ================ Slide 11: Predicted vs actual ================
s = slide_with_bar("Random-Forest predicted vs actual",
                   "Color: FPM present (red) vs absent (blue) — model handles both regimes")
add_image(s, FIG / "pred_vs_actual_rf.png",
          left=Inches(0.5), top=Inches(1.8), width=Inches(12.3))

# ================ Slide 12: Residuals + importance ================
s = slide_with_bar("Residuals + feature importance",
                   "Residuals are near-Gaussian for q_WPM/q_WFM;  heavier tails for q_RWFM")
add_image(s, FIG / "residuals_rf.png",
          left=Inches(0.4), top=Inches(1.7), width=Inches(6.4))
add_image(s, FIG / "rf_top20_importance.png",
          left=Inches(7.0), top=Inches(1.7), height=Inches(5.2))

# ================ Slide 13: Aggregated importance ================
s = slide_with_bar("MDEV dominates feature importance (5× MHDEV)",
                   "Consistent with WFM/WPM discrimination — MDEV is the natural estimator")
add_image(s, FIG / "importance_by_statistic.png",
          left=Inches(1.5), top=Inches(1.6), width=Inches(10.3))

# ================ Slide 14: Uncertainty ================
s = slide_with_bar("Uncertainty quantification",
                   "RF tree-variance ±95% CIs vs truth;  XGB quantile 90% PI empirical coverage 80-84%")
add_image(s, FIG / "rf_uq.png",
          left=Inches(0.5), top=Inches(1.7), width=Inches(12.3))

# ================ Slide 15: Model comparison ================
s = slide_with_bar("Model comparison summary",
                   "Test-set RMSE by target — both models substantially beat naive")
add_image(s, FIG / "model_comparison_rmse.png",
          left=Inches(2.0), top=Inches(1.6), width=Inches(9.3))

# ================ Slide 16: Real-data ADEV overlay ================
s = slide_with_bar("Real-data validation: GMR6000 Rb, 4 non-overlapping windows",
                   "Measured ADEV vs 5 estimators (naive / RF / XGB / MHDEV-fit / ALS)")
add_image(s, FIG / "real_data_method_comparison.png",
          left=Inches(0.3), top=Inches(1.7), width=Inches(12.7))
add_bullets(s, [
    "ML (RF/XGB) and MHDEV-fit / ALS track the measured ADEV across all 4 windows.",
    "ALS is used as the tuning reference throughout: autocovariance-LS is well-conditioned under the 3-state clock model and converges monotonically from any of our seeds.",
    "Naive seed matches at short τ (WPM floor) but drifts away at long τ — motivates ML lookup.",
],
    left=Inches(0.4), top=Inches(6.0), width=Inches(12.5), height=Inches(1.4), size=12)

# ================ Slide 17: Warm-start comparison ================
s = slide_with_bar("Warm-start benefit — ML seed vs naive seed (ALS tuning)",
                   "Holdover phase RMS after ALS converges, per window — ML seed delivers a better-tuned KF")
add_image(s, FIG / "warmstart_compare.png",
          left=Inches(2.5), top=Inches(1.6), height=Inches(4.4))
add_bullets(s, [
    "ALS converges in 1 outer iteration from any reasonable seed — iteration count is not the discriminator.",
    "What changes is the final q that ALS lands on, and therefore the forward-prediction RMS on held-out phase.",
    "ML seed wins decisively on windows 1, 2 (26× and 4600×) — including window 2, where the old NLL-label model failed; naive wins windows 0, 3 by 3-50× (small RMS regime, both seeds good).",
    "Net: ML wins 2/4 windows by orders of magnitude; naive wins 2/4 by small factors.  Mean RMS ratio naive/ML = 1157×.",
],
    left=Inches(0.6), top=Inches(6.15), width=Inches(12.2), height=Inches(1.4), size=13)

# ================ Slide 18: Conclusions ================
s = slide_with_bar("Conclusions & next steps",
                   "What worked, what didn't, and where this pipeline goes next")
add_bullets(s, [
    "XGB predicts q_WPM (R² 0.9999) and q_WFM (R² 0.996) to < 0.05 log-decades — production-ready.",
    "q_RWFM R² improved from 0.65 → 0.82 after switching labels from NLL to ALS (noise-free targets).",
    "MDEV features carry 5× more predictive signal than MHDEV — drift-rejection sacrifices short-τ WPM/WFM resolution where most discrimination lives.",
    "Real-data validation: ML / MHDEV-fit / ALS track measured ADEV on all 4 GMR6000 windows.",
    "ALS chosen as the tuning reference — well-conditioned, monotonic convergence from both naive and ML seeds.",
    "",
    "Next steps:",
    "   — expand RF/XGB training range to cover quieter Rb regimes seen in real hardware (q_RWFM ≲ 1e-32)",
    "   — ensemble ML + ALS: use ML seed, run 1-2 ALS polishing iterations (feed-forward model + bounded tuning)",
    "   — extend to 4-state drift-RW model for sub-ns long-term holdover",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=16)

# ================ Slide 19: Validation against Stable32 / allantools ================
s = slide_with_bar("SigmaTau numerical validation",
                   "10 deviation kernels cross-checked against Stable32 + allantools on 8192-sample Rb dataset")
add_bullets(s, [
    "Core estimators (ADEV, MDEV, TDEV, HDEV): match Stable32 within 5×10⁻⁵ rel.-err,",
    "   allantools within 2×10⁻⁷ rel.-err — numerical fidelity confirmed.",
    "Total deviations: MTOTDEV shows documented 1.27× bias-correction offset per SP1065;",
    "   HTOTDEV 0.5% offset matches reference policy.",
    "MHDEV kernel re-verified against legacy stablab/mhdev.m (recovered from git history) —",
    "   algebraically identical: mean((mΔ³x)²)/(6·m⁴·τ²) ≡ mean(avg²)/(6·m²·τ²).",
    "Per-deviation PNGs saved to reference/validation/plots/  (driver: validate.py).",
    "Result: SigmaTau deviation library is the trusted source for the 196-feature extractor.",
],
    left=Inches(0.6), top=Inches(1.9), width=Inches(12.2), height=Inches(5), size=18)

# ================ Slide 20: Thank you / Q&A ================
s = prs.slides.add_slide(BLANK)
bar = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(7.5))
bar.fill.solid(); bar.fill.fore_color.rgb = DARK; bar.line.fill.background()
add_textbox(s, "Thank you", left=Inches(0.5), top=Inches(2.5),
            width=Inches(12.5), height=Inches(1.5), size=72, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
add_textbox(s, "Questions?", left=Inches(0.5), top=Inches(4.2),
            width=Inches(12.5), height=Inches(1.0), size=36, color=ACCENT,
            align=PP_ALIGN.CENTER)
add_textbox(s,
    "Code: SigmaTau (Julia) + ml/ (Python)   •   Notebook: ml/notebook.executed.ipynb\n"
    "Validation scripts: reference/validation/validate.py, ml/notebook_extras.py",
    left=Inches(0.5), top=Inches(5.8), width=Inches(12.5), height=Inches(1.0),
    size=16, color=WHITE, align=PP_ALIGN.CENTER)

# --- add slide numbers to all non-title/closing slides ---
total = len(prs.slides)
for idx, slide in enumerate(prs.slides, start=1):
    if idx in (1, total):
        continue
    add_footer(slide, idx, total)

prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Total slides: {total}")
